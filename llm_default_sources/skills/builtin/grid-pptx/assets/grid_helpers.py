"""
GRID PPTX helpers — a Python module for generating GRID-branded presentations.

Usage from Claude's build script:
    import sys
    sys.path.insert(0, "/path/to/skill/assets")
    import grid_helpers as g

    g.init()  # uses the bundled template.pptx next to this file
    # OR: g.init("/custom/path/to/template.pptx")

    g.remove_template_slide()
    g.add_title_slide("My Deck", "Subtitle")
    g.add_content_slide("Agenda", ["Intro", "Body", "Outro"])
    g.save("/path/to/output.pptx")

All layout/helper/style functions are exposed as module-level names after init().
"""

from pathlib import Path

# ═══════════════════════════════════════════════════════════
#  MODULE STATE (populated by init())
# ═══════════════════════════════════════════════════════════
prs = None
_template_path = None
L_TITLE = L_SECTION = L_TEXT = L_BLANK = L_SCREEN = None

# These get set on import so callers can reference them before init()
Presentation = None
Inches = Pt = Emu = None
RGBColor = None
PP_ALIGN = MSO_ANCHOR = None
MSO_SHAPE = None


def _load_pptx_lib():
    """Import python-pptx. Assumes it's already installed (SKILL.md handles install)."""
    global Presentation, Inches, Pt, Emu, RGBColor, PP_ALIGN, MSO_ANCHOR, MSO_SHAPE
    from pptx import Presentation as _Presentation
    from pptx.util import Inches as _Inches, Pt as _Pt, Emu as _Emu
    from pptx.dml.color import RGBColor as _RGBColor
    from pptx.enum.text import PP_ALIGN as _PP_ALIGN, MSO_ANCHOR as _MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE as _MSO_SHAPE

    Presentation = _Presentation
    Inches, Pt, Emu = _Inches, _Pt, _Emu
    RGBColor = _RGBColor
    PP_ALIGN, MSO_ANCHOR = _PP_ALIGN, _MSO_ANCHOR
    MSO_SHAPE = _MSO_SHAPE


def init(template_path=None):
    """Initialize the presentation from the GRID template.

    Args:
        template_path: Path to template.pptx. If None, uses the bundled template.pptx
                       located next to this module file.

    Returns:
        The pptx.Presentation object.
    """
    global prs, _template_path, L_TITLE, L_SECTION, L_TEXT, L_BLANK, L_SCREEN
    global C, F  # color/font classes (created after RGBColor is imported)

    _load_pptx_lib()

    if template_path is None:
        template_path = str(Path(__file__).parent / "template.pptx")
    _template_path = template_path

    if not Path(template_path).exists():
        raise FileNotFoundError(
            f"GRID template not found: {template_path}\n"
            "Please provide the template.pptx file or use the bundled one."
        )

    prs = Presentation(template_path)

    # Layout shortcuts
    L_TITLE   = prs.slide_layouts[0]   # Title (cover) — dark bg
    L_SECTION = prs.slide_layouts[1]   # Section Title — dark bg
    L_TEXT    = prs.slide_layouts[2]   # Contents Text — white bg, title + body
    L_BLANK   = prs.slide_layouts[3]   # Contents Blank — white bg, title only
    L_SCREEN  = prs.slide_layouts[4]   # Screen Shot — white bg, fully blank

    # Build color and font classes now that RGBColor is loaded
    _build_brand_classes()

    return prs


def _build_brand_classes():
    """Populate module-level C (colors) and F (fonts) classes."""
    global C, F

    class _C:
        """GRID2026 theme colors. Usage: g.C.blue, g.C.navy, etc."""
        navy       = RGBColor(0x1D, 0x27, 0x49)  # dk2  — dark backgrounds, table headers
        dark_navy  = RGBColor(0x1A, 0x27, 0x4F)  # layout background fill
        blue       = RGBColor(0x2F, 0x6E, 0xBA)  # accent1 — primary blue, headers
        cyan       = RGBColor(0x30, 0x93, 0xC6)  # accent2 — secondary blue
        green      = RGBColor(0x05, 0xAF, 0x50)  # accent3 — positive/success
        gold       = RGBColor(0xFF, 0xC0, 0x00)  # accent4 — warnings, highlights
        orange     = RGBColor(0xDF, 0x76, 0x2E)  # accent5 — emphasis
        red        = RGBColor(0xF5, 0x5F, 0x1D)  # accent6 — alerts
        text       = RGBColor(0x59, 0x59, 0x59)  # dk1  — body text on light bg
        gray       = RGBColor(0xAC, 0xAC, 0xAC)  # lt2  — muted/copyright text
        white      = RGBColor(0xFF, 0xFF, 0xFF)  # lt1
        black      = RGBColor(0x00, 0x00, 0x00)
        link       = RGBColor(0x05, 0x63, 0xC1)  # hyperlink blue
        row_alt    = RGBColor(0xF2, 0xF2, 0xF2)  # light gray for zebra striping

    class _F:
        """GRID brand fonts."""
        heading = "Yu Gothic"    # Major / headings (游ゴシック)
        body    = "Calibri"      # Minor / body text
        shape   = "Meiryo UI"    # Default shape text
        mono    = "Consolas"     # Code / technical text

    C = _C
    F = _F


# Declare C and F up front so IDEs/Claude can reference them
C = None
F = None


# ═══════════════════════════════════════════════════════════
#  STYLE HELPERS
# ═══════════════════════════════════════════════════════════

def style(run, size=None, bold=None, italic=None, color=None, font=None):
    """Apply formatting to a single text run.

    Args:
        run:    pptx Run object
        size:   font size in pt (int)
        bold:   True/False/None
        italic: True/False/None
        color:  RGBColor object (use C.blue, C.text, etc.)
        font:   font family name (str) — sets latin, ea, and cs fonts
    """
    if size is not None:
        run.font.size = Pt(size)
    if bold is not None:
        run.font.bold = bold
    if italic is not None:
        run.font.italic = italic
    if color is not None:
        run.font.color.rgb = color
    if font is not None:
        run.font.name = font
        # Also set East Asian and Complex Script fonts for JP support
        from pptx.oxml.ns import qn
        from lxml import etree
        rPr = run._r.get_or_add_rPr()
        for tag in ['a:ea', 'a:cs']:
            el = rPr.find(qn(tag))
            if el is None:
                el = etree.SubElement(rPr, qn(tag))
            el.set('typeface', font)


def add_paragraph(tf, text, level=0, size=14, bold=False, italic=False,
                  color=None, font=None, alignment=None):
    """Add a new paragraph to a text frame. Returns the paragraph."""
    p = tf.add_paragraph()
    p.text = text
    p.level = level
    if alignment:
        p.alignment = alignment
    if p.runs:
        style(p.runs[0], size=size, bold=bold, italic=italic, color=color, font=font)
    return p


def set_paragraph(tf, text, size=14, bold=False, color=None, font=None, alignment=None):
    """Replace text frame content with a single paragraph. Returns the paragraph."""
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    if alignment:
        p.alignment = alignment
    if p.runs:
        style(p.runs[0], size=size, bold=bold, color=color, font=font)
    return p


# ═══════════════════════════════════════════════════════════
#  SLIDE BUILDERS
# ═══════════════════════════════════════════════════════════

def add_title_slide(title, subtitle=""):
    """Cover slide — dark navy background with wireframe sphere.
    Use for the very first slide of a presentation.
    """
    slide = prs.slides.add_slide(L_TITLE)
    slide.placeholders[0].text = title
    if subtitle:
        slide.placeholders[1].text = subtitle
    return slide


def add_section_slide(title, subtitle=""):
    """Section divider — dark navy background with wireframe sphere.
    Use to separate major sections of the presentation.
    """
    slide = prs.slides.add_slide(L_SECTION)
    slide.placeholders[0].text = title
    if subtitle:
        slide.placeholders[1].text = subtitle
    return slide


def add_content_slide(title, bullets=None):
    """Content slide with title + bullet list on white background.

    Args:
        title:   slide title string
        bullets: list of items, each item is either:
                 - str              → level-0 bullet
                 - (str, int)       → (text, indent_level)
                 - (str, int, dict) → (text, indent_level, {size, bold, color, font})
    Returns: slide object
    """
    slide = prs.slides.add_slide(L_TEXT)
    slide.placeholders[0].text = title
    if bullets:
        tf = slide.placeholders[1].text_frame
        tf.clear()
        for i, item in enumerate(bullets):
            # Parse item
            if isinstance(item, str):
                text, level, opts = item, 0, {}
            elif len(item) == 2:
                text, level = item
                opts = {}
            else:
                text, level, opts = item[0], item[1], item[2]

            if i == 0:
                # Use existing first paragraph
                p = tf.paragraphs[0]
                p.text = text
                p.level = level
                if p.runs:
                    style(p.runs[0], size=opts.get("size", 14),
                          bold=opts.get("bold", False),
                          color=opts.get("color"), font=opts.get("font"))
            else:
                add_paragraph(tf, text, level=level,
                              size=opts.get("size", 14),
                              bold=opts.get("bold", False),
                              color=opts.get("color"), font=opts.get("font"))
    return slide


def add_blank_slide(title=""):
    """Blank canvas slide — white bg, title only, no body placeholder.
    Use for custom layouts with shapes, images, charts, etc.
    """
    slide = prs.slides.add_slide(L_BLANK)
    if title:
        slide.placeholders[0].text = title
    return slide


def add_screenshot_slide():
    """Fully blank slide — no title, no placeholders.
    Use for full-bleed images or screenshots.
    """
    return prs.slides.add_slide(L_SCREEN)


def add_table_slide(title, headers, rows, col_widths=None, zebra=True):
    """Slide with a professionally styled table (navy header, zebra striping).

    Args:
        title:      slide title string
        headers:    list of column header strings
        rows:       list of lists — each inner list is one row of cell values
        col_widths: optional list of floats (inches) per column
        zebra:      alternate row background colors (default True)
    Returns: slide object
    """
    slide = prs.slides.add_slide(L_BLANK)
    slide.placeholders[0].text = title

    n_rows = len(rows) + 1
    n_cols = len(headers)
    table_height = min(n_rows * 0.6, 5.5)

    tbl_shape = slide.shapes.add_table(
        n_rows, n_cols,
        Inches(0.5), Inches(1.0), Inches(12.3), Inches(table_height)
    )
    table = tbl_shape.table

    if col_widths:
        for j, w in enumerate(col_widths):
            table.columns[j].width = Inches(w)

    # Header row
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = C.navy
        for p in cell.text_frame.paragraphs:
            for run in p.runs:
                style(run, size=13, bold=True, color=C.white, font=F.body)

    # Data rows
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.text = str(val)
            if zebra and i % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = C.row_alt
            for p in cell.text_frame.paragraphs:
                for run in p.runs:
                    style(run, size=12, color=C.text, font=F.body)
    return slide


# ═══════════════════════════════════════════════════════════
#  SHAPE & TEXT BOX HELPERS
# ═══════════════════════════════════════════════════════════

def add_textbox(slide, text, x, y, w, h, size=14, bold=False, color=None,
                font=None, alignment=None, word_wrap=True):
    """Add a text box at exact position. Returns the shape.

    Positions in inches. Defaults: alignment=PP_ALIGN.LEFT, color=C.text, font=F.body.
    """
    if alignment is None:
        alignment = PP_ALIGN.LEFT
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = alignment
    if p.runs:
        style(p.runs[0], size=size, bold=bold,
              color=color or C.text, font=font or F.body)
    return txBox


def add_shape(slide, text, x, y, w, h, fill_color=None, text_color=None,
              font_size=16, font=None, shape_type=None,
              alignment=None, anchor=None):
    """Add a shape with centered text. Returns the shape.

    Positions in inches. Defaults: shape_type=RECTANGLE, alignment=CENTER, anchor=MIDDLE.
    """
    if shape_type is None:
        shape_type = MSO_SHAPE.RECTANGLE
    if alignment is None:
        alignment = PP_ALIGN.CENTER
    if anchor is None:
        anchor = MSO_ANCHOR.MIDDLE

    shape = slide.shapes.add_shape(
        shape_type, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    shape.line.fill.background()

    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    tf._txBody.bodyPr.set('anchor', {
        MSO_ANCHOR.TOP: 't', MSO_ANCHOR.MIDDLE: 'ctr', MSO_ANCHOR.BOTTOM: 'b'
    }.get(anchor, 'ctr'))

    p = tf.paragraphs[0]
    p.text = text
    p.alignment = alignment
    if p.runs:
        style(p.runs[0], size=font_size,
              color=text_color or C.white, font=font or F.body)
    return shape


def add_image(slide, image_path, x, y, w=None, h=None):
    """Add an image to a slide. Returns the shape.
    Specify w or h (or both). If only one is given, aspect ratio is preserved.
    """
    kwargs = {}
    if w: kwargs['width'] = Inches(w)
    if h: kwargs['height'] = Inches(h)
    return slide.shapes.add_picture(image_path, Inches(x), Inches(y), **kwargs)


# ═══════════════════════════════════════════════════════════
#  COMPOSITE CARD HELPERS (useful for conversion workflows)
# ═══════════════════════════════════════════════════════════

def add_card_with_top_border(slide, x, y, w, h, border_color, border_h=0.06):
    """Card with light background and thin colored top border.
    Returns the card shape (not the border).
    """
    card = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(0xE8, 0xED, 0xF5)
    card.line.fill.background()
    border = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(border_h)
    )
    border.fill.solid()
    border.fill.fore_color.rgb = border_color
    border.line.fill.background()
    return card


def add_card_with_left_bar(slide, x, y, w, h, bar_color, bar_w=0.10):
    """Card with light background and thin colored left bar.
    Returns the card shape (not the bar).
    """
    card = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(0xE8, 0xED, 0xF5)
    card.line.fill.background()
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(bar_w), Inches(h)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = bar_color
    bar.line.fill.background()
    return card


def add_text_arrow(slide, char, x, y, w, h, size=14, color=None):
    """Add a ▶ or ▼ text character as an arrow."""
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = char
    p.alignment = PP_ALIGN.CENTER
    if p.runs:
        r = p.runs[0]
        r.font.size = Pt(size)
        r.font.color.rgb = color or RGBColor(0x64, 0x74, 0x8B)
    return txBox


# ═══════════════════════════════════════════════════════════
#  TWO-COLUMN & MULTI-BOX HELPERS
# ═══════════════════════════════════════════════════════════

def add_two_column_slide(title, left_items, right_items, left_header="", right_header=""):
    """Content slide with two columns. Uses Contents Blank layout."""
    slide = add_blank_slide(title)
    col_w = 5.8
    left_x, right_x = 0.5, 6.8
    start_y = 1.0

    for items, x, header in [
        (left_items, left_x, left_header),
        (right_items, right_x, right_header),
    ]:
        y_offset = start_y
        if header:
            add_textbox(slide, header, x, y_offset, col_w, 0.4,
                        size=16, bold=True, color=C.blue, font=F.heading)
            y_offset += 0.5
        for item in items:
            add_textbox(slide, f"• {item}", x + 0.2, y_offset, col_w - 0.2, 0.35,
                        size=13, color=C.text)
            y_offset += 0.4
    return slide


def add_kpi_slide(title, kpis):
    """Slide with large KPI callout numbers. Uses Contents Blank layout.

    Args:
        kpis: list of dicts: [{"value": "25%", "label": "Growth", "color": C.green}, ...]
              Max 4 KPIs recommended.
    """
    slide = add_blank_slide(title)
    n = len(kpis)
    box_w = min(3.0, 12.0 / n - 0.3)
    total_w = n * box_w + (n - 1) * 0.3
    start_x = (13.33 - total_w) / 2

    for i, kpi in enumerate(kpis):
        x = start_x + i * (box_w + 0.3)
        color = kpi.get("color", C.blue)
        add_shape(slide, kpi["value"], x, 2.0, box_w, 2.5,
                  fill_color=color, text_color=C.white, font_size=48, font=F.heading)
        add_textbox(slide, kpi["label"], x, 4.7, box_w, 0.5,
                    size=14, bold=True, color=C.text, alignment=PP_ALIGN.CENTER)
    return slide


# ═══════════════════════════════════════════════════════════
#  SAVE & CLEANUP
# ═══════════════════════════════════════════════════════════

def remove_template_slide():
    """Remove the template's default blank first slide.
    Call BEFORE adding your slides to avoid shipping a blank placeholder at the top.
    """
    if len(prs.slides._sldIdLst) > 0:
        from pptx.oxml.ns import qn
        sldId = prs.slides._sldIdLst[0]
        rId = sldId.get(qn('r:id'))
        prs.part.drop_rel(rId)
        prs.slides._sldIdLst.remove(sldId)


def update_copyright_year(new_year):
    """Replace '© 2022 GRID Inc.' in the template with the current year.
    Call BEFORE adding any slides."""
    target = "© 2022"
    replacement = f"© {new_year}"
    # Scan layouts and masters
    sources = list(prs.slide_layouts) + list(prs.slide_masters)
    for src in sources:
        for shape in src.shapes:
            if shape.has_text_frame:
                for p in shape.text_frame.paragraphs:
                    for run in p.runs:
                        if target in (run.text or ""):
                            run.text = run.text.replace(target, replacement)


def save(output_path):
    """Save the presentation to disk."""
    prs.save(output_path)
    print(f"Saved: {output_path}")
    return output_path
